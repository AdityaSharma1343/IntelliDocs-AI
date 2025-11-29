"""
Document Processor for various file formats
"""

import os
import logging
from typing import Dict, Any

# Set up logger
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Import document libraries with error handling
try:
    import PyPDF2
    from docx import Document
    from pptx import Presentation
    import openpyxl
    from openpyxl import load_workbook
    logger.info("✅ All document processors imported successfully")
except ImportError as e:
    logger.warning(f"⚠️ Some document processors missing: {e}")


class DocumentProcessor:
    """Process different document types and extract text"""
    
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """Extract text from PDF file"""
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                num_pages = len(pdf_reader.pages)
                
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            
            return text.strip() or "No text content extracted from PDF"
        except Exception as e:
            logger.error(f"Error processing PDF: {e}")
            return f"Error extracting text from PDF: {str(e)}"
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """Extract text from Word document"""
        try:
            doc = Document(file_path)
            text = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text)
            
            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text.append(cell.text)
            
            return "\n".join(text) or "No text content extracted from Word document"
        except Exception as e:
            logger.error(f"Error processing DOCX: {e}")
            return f"Error extracting text from Word: {str(e)}"
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """Extract text from PowerPoint presentation"""
        try:
            prs = Presentation(file_path)
            text = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                text.append(f"\n--- Slide {slide_num} ---\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text.append(shape.text)
            
            return "\n".join(text) or "No content extracted from PowerPoint"
        except Exception as e:
            logger.error(f"Error processing PPTX: {e}")
            return f"Error extracting text from PowerPoint: {str(e)}"
    
    @staticmethod
    def extract_text_from_excel(file_path: str) -> str:
        """Extract text from Excel file"""
        try:
            # Import here as backup
            try:
                import openpyxl
                from openpyxl import load_workbook
            except ImportError:
                return "❌ Error: openpyxl library not installed"
            
            workbook = load_workbook(file_path, data_only=True, read_only=True)
            text = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text.append(f"\n{'='*50}")
                text.append(f"Sheet: {sheet_name}")
                text.append(f"{'='*50}\n")
                
                row_count = 0
                for row in sheet.iter_rows(values_only=True):
                    row_values = []
                    for cell in row:
                        if cell is not None:
                            row_values.append(str(cell))
                        else:
                            row_values.append("")
                    
                    if any(val.strip() for val in row_values):
                        row_text = " | ".join(row_values)
                        text.append(row_text)
                        row_count += 1
                
                text.append(f"\n(Total rows: {row_count})\n")
            
            workbook.close()
            
            result = "\n".join(text)
            return result if result.strip() else "No content extracted from Excel"
            
        except ImportError as e:
            error_msg = f"❌ Import Error: openpyxl not available - {str(e)}"
            logger.error(error_msg)
            return error_msg
        except Exception as e:
            logger.error(f"Error processing Excel: {e}")
            return f"Error extracting text from Excel: {str(e)}"
    
    @staticmethod
    def extract_text_from_txt(file_path: str) -> str:
        """Extract text from plain text file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            logger.error(f"Error processing TXT: {e}")
            return f"Error reading text file: {str(e)}"
    
    @staticmethod
    def extract_text_from_csv(file_path: str) -> str:
        """Extract text from CSV file"""
        try:
            import csv
            text = []
            with open(file_path, 'r', encoding='utf-8', newline='') as csvfile:
                reader = csv.reader(csvfile)
                for row_num, row in enumerate(reader, 1):
                    row_text = " | ".join(str(cell) for cell in row if cell)
                    if row_text.strip():
                        text.append(f"Row {row_num}: {row_text}")
            return "\n".join(text) or "No content extracted from CSV"
        except Exception as e:
            logger.error(f"Error processing CSV: {e}")
            return f"Error extracting text from CSV: {str(e)}"

    @staticmethod
    def process_document(file_path: str) -> Dict[str, Any]:
        """
        Process any document type and extract text
        Returns dict with text and metadata
        """
        file_extension = os.path.splitext(file_path)[1].lower()
        filename = os.path.basename(file_path)
        
        logger.info(f"Processing document: {filename} (type: {file_extension})")
        
        # Extract text based on file type
        text = ""
        doc_type = "unknown"
        
        if file_extension == '.pdf':
            text = DocumentProcessor.extract_text_from_pdf(file_path)
            doc_type = "PDF"
        elif file_extension in ['.docx', '.doc']:
            text = DocumentProcessor.extract_text_from_docx(file_path)
            doc_type = "Word"
        elif file_extension in ['.pptx', '.ppt']:
            text = DocumentProcessor.extract_text_from_pptx(file_path)
            doc_type = "PowerPoint"
        elif file_extension in ['.xlsx', '.xls']:
            text = DocumentProcessor.extract_text_from_excel(file_path)
            doc_type = "Excel"
        elif file_extension == '.txt':
            text = DocumentProcessor.extract_text_from_txt(file_path)
            doc_type = "Text"
        elif file_extension == '.csv':
            text = DocumentProcessor.extract_text_from_csv(file_path)
            doc_type = "CSV"
        else:
            text = f"Unsupported file type: {file_extension}"
            doc_type = "Unsupported"
        
        # Get file metadata
        file_size = os.path.getsize(file_path) if os.path.exists(file_path) else 0
        
        logger.info(f"Extracted {len(text)} characters from {filename}")
        
        return {
            "filename": filename,
            "content": text,
            "doc_type": doc_type,
            "file_size": file_size,
            "word_count": len(text.split()) if text else 0,
            "char_count": len(text)
        }


# Test function
if __name__ == "__main__":
    # Test with a sample file
    test_file = "sample.pdf"  # Change to your test file
    if os.path.exists(test_file):
        result = DocumentProcessor.process_document(test_file)
        print(f"Document: {result['filename']}")
        print(f"Type: {result['doc_type']}")
        print(f"Words: {result['word_count']}")
        print(f"Preview: {result['content'][:200]}...")
    else:
        print(f"Test file not found: {test_file}")